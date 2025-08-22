from hdbcli import dbapi
import csv
from datetime import date
from timestamp import add_timestamp_to_filename 
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from scipy.stats import chisquare, chi2, norm
from sklearn.mixture import GaussianMixture
import matplotlib.dates as mdates
from sklearn.linear_model import LinearRegression
import seaborn as sns
import json
import os
from pptx import Presentation
from pptx.util import Inches
import os
import re


def execute_query(connection, query, params=None):
    cursor = connection.cursor()
    cursor.execute(query, params or ())
    results = cursor.fetchall()
    #column_names = [desc[0] for desc in cursor.description]
    #print("Spaltennamen:", column_names)
    cursor.close()
    return results




def bestandsverlauf_pro_LO_tag(df):
    """
    Erstellt einen vollständigen Bestandsverlauf pro Tag und Lagerort.
    Falls für einen Tag kein Eintrag existiert, wird der letzte bekannte Wert übernommen.
    Falls kein vorheriger Wert existiert, wird 0 gesetzt.

    :param df: Pandas DataFrame mit 'MVMT_STORAGE_AREA', 'MVMT_STORAGE_AREA_ON_HAND', 'MVMT_CREATION_DATE_TIME'
    :return: DataFrame mit vollständigem Bestandsverlauf pro Tag
    """

    # **1. Datumsformat umwandeln**
    df['MVMT_CREATION_DATE_TIME'] = pd.to_datetime(df['MVMT_CREATION_DATE_TIME'], format='%Y-%m-%d %H:%M:%S', errors='coerce')

    # **2. Neues Feld für das Datum ohne Uhrzeit**
    df['DATE'] = df['MVMT_CREATION_DATE_TIME'].dt.date

    # **3. Daten nach Lagerort und Datum sortieren**
    df_sorted = df.sort_values(by=['MVMT_STORAGE_AREA', 'DATE', 'MVMT_CREATION_DATE_TIME'])

    # **4. Letzten Eintrag pro Tag & Lagerort nehmen**
    df_daily_stock = df_sorted.groupby(['MVMT_STORAGE_AREA', 'DATE']).last().reset_index()

    # **5. Vollständige Datumsreihe erzeugen**
    min_date = df_daily_stock['DATE'].min()
    max_date = df_daily_stock['DATE'].max()
    all_dates = pd.date_range(start=min_date, end=max_date).date

    # **6. Fehlende Daten für jeden Lagerort auffüllen**
    all_storage_areas = df_daily_stock['MVMT_STORAGE_AREA'].unique()
    complete_data = []

    for storage_area in all_storage_areas:
        subset = df_daily_stock[df_daily_stock['MVMT_STORAGE_AREA'] == storage_area].set_index('DATE')
        subset = subset.reindex(all_dates)  # Fehltage auffüllen
        subset['MVMT_STORAGE_AREA'] = storage_area  # Lagerort ergänzen
        subset['MVMT_STORAGE_AREA_ON_HAND'].ffill(inplace=True)  # Fehlende Werte mit vorherigem Bestand füllen
        subset['MVMT_STORAGE_AREA_ON_HAND'].fillna(0, inplace=True)  # Falls es noch keinen vorherigen Wert gab → 0 setzen
        subset.reset_index(inplace=True)  # Index zurück in eine Spalte umwandeln und nicht als Index lassen
        subset.rename(columns={'index': 'DATE'}, inplace=True)  # Datumsspalte korrekt benennen
        complete_data.append(subset)

    # **7. Alles in ein neues DataFrame packen**
    df_complete = pd.concat(complete_data, ignore_index=True)

    # **8. Diagramm für Bestandsverlauf pro Lagerort plotten**
    plt.figure(figsize=(12, 6))
    
    for storage_area in df_complete['MVMT_STORAGE_AREA'].unique():
        subset = df_complete[df_complete['MVMT_STORAGE_AREA'] == storage_area]
        plt.plot(subset['DATE'], subset['MVMT_STORAGE_AREA_ON_HAND'], marker='o', linestyle='-', label=f'Lager {storage_area}')

    # **9. Diagramm-Formatierung**
    plt.xlabel('Datum')
    plt.ylabel('Bestand')
    plt.title('Bestandsverlauf pro Lagerort mit vollständiger Zeitreihe')
    plt.xticks(rotation=45)
    plt.legend()
    plt.grid(True)

    # **10. Diagramm anzeigen**
    #plt.show()

    return df_complete

def gesamter_bestandsverlauf(df):
    """
    Erstellt einen vollständigen Bestandsverlauf pro Tag basierend auf 'MVMT_ON_HAND'.
    Falls für einen Tag kein Eintrag existiert, wird der letzte bekannte Wert übernommen.
    Falls kein vorheriger Wert existiert, wird 0 gesetzt.

    :param df: Pandas DataFrame mit 'MVMT_ON_HAND' und 'MVMT_CREATION_DATE_TIME'
    :return: DataFrame mit vollständigem Bestandsverlauf pro Tag
    """

    # **1. Datumsformat umwandeln**
    df['MVMT_CREATION_DATE_TIME'] = pd.to_datetime(df['MVMT_CREATION_DATE_TIME'], format='%Y-%m-%d %H:%M:%S', errors='coerce')

    # **2. Neues Feld für das Datum ohne Uhrzeit**
    df['DATE'] = df['MVMT_CREATION_DATE_TIME'].dt.date

    # **3. Daten nach Datum und Zeit sortieren**
    df_sorted = df.sort_values(by=['DATE', 'MVMT_CREATION_DATE_TIME'])

    # **4. Letzten Bestand pro Tag bestimmen (aus 'MVMT_ON_HAND')**
    df_daily_stock = df_sorted.groupby('DATE')['MVMT_ON_HAND'].last().reset_index()

    # **5. Vollständige Datumsreihe erzeugen**
    min_date = df_daily_stock['DATE'].min()
    max_date = df_daily_stock['DATE'].max()
    all_dates = pd.date_range(start=min_date, end=max_date).date

    # **6. Fehlende Daten auffüllen**
    df_daily_stock = df_daily_stock.set_index('DATE').reindex(all_dates)  # Fehlende Tage hinzufügen
    df_daily_stock['MVMT_ON_HAND'].ffill(inplace=True)  # Fehlende Werte mit vorherigem Bestand füllen
    df_daily_stock['MVMT_ON_HAND'].fillna(0, inplace=True)  # Falls es noch keinen vorherigen Wert gab → 0 setzen
    df_daily_stock.reset_index(inplace=True)  # Index zurück in eine Spalte umwandeln
    df_daily_stock.rename(columns={'index': 'DATE'}, inplace=True)  # Spalte korrekt benennen

    # **7. Diagramm für Gesamtbestandsverlauf plotten**
    plt.figure(figsize=(12, 6))
    plt.plot(df_daily_stock['DATE'], df_daily_stock['MVMT_ON_HAND'], marker='o', linestyle='-', color='b', label='Gesamtbestand')

    # **8. Diagramm-Formatierung**
    plt.xlabel('Datum')
    plt.ylabel('Gesamtbestand')
    plt.title('Gesamtbestandsverlauf basierend auf MVMT_ON_HAND')
    plt.xticks(rotation=45)
    plt.legend()
    plt.grid(True)

    # **9. Diagramm anzeigen**
    #plt.show()

    return df_daily_stock

# **Beispielaufruf**
# df_gesamtbestand = gesamter_bestandsverlauf(df)


def bestandsverlauf_moving_average(df, window=14):
    """
    Erstellt einen Bestandsverlauf mit gleitendem Mittelwert über eine bestimmte Anzahl an Tagen.

    :param df: Pandas DataFrame mit 'DATE' und 'MVMT_ON_HAND'
    :param window: Fenstergröße für den gleitenden Mittelwert (Standard: 14 Tage)
    :return: DataFrame mit gleitendem Mittelwert
    """

    # **1. Sicherstellen, dass DATE eine Zeitreihe ist**
    df['DATE'] = pd.to_datetime(df['DATE'])

    # **2. Gleitenden Mittelwert berechnen (14-Tage-Schnitt)**
    df['MVMT_ON_HAND_MA'] = df['MVMT_ON_HAND'].rolling(window=window, min_periods=1).mean()

    # **3. Diagramm für Bestandsverlauf mit gleitendem Mittelwert plotten**
    plt.figure(figsize=(12, 6))

    # Original-Bestand plotten
    plt.plot(df['DATE'], df['MVMT_ON_HAND'], marker='o', linestyle='-', color='b', alpha=0.4, label='Täglicher Bestand')

    # Gleitenden Mittelwert plotten
    plt.plot(df['DATE'], df['MVMT_ON_HAND_MA'], marker='', linestyle='-', color='r', linewidth=2, label=f'{window}-Tage Moving Average')

    # **4. Diagramm-Formatierung**
    plt.xlabel('Datum')
    plt.ylabel('Bestand')
    plt.title(f'Bestandsverlauf mit {window}-Tage Moving Average')
    plt.xticks(rotation=45)
    plt.legend()
    plt.grid(True)

    # **5. Diagramm anzeigen**
    #plt.show()

    return df

def plot_scatter_with_regression(df, titletext, part, json_file="regression_results.json"):
    """
    Erstellt einen Scatter-Plot mit linearer Regression und speichert die Steigung (slope) in einer JSON-Datei.
    
    :param df: Pandas DataFrame mit 'CUM_QTY' und 'REAL_QTY'
    :param titletext: Titel für den Plot
    :param part: ID des Teils (wird als Key in der JSON-Datei genutzt)
    :param json_file: Name der JSON-Datei zur Speicherung der Regressionsergebnisse
    """

    # ✅ Sicherstellen, dass die relevanten Spalten existieren
    if 'CUM_QTY' not in df.columns or 'REAL_QTY' not in df.columns:
        raise ValueError("Fehlende Spalten: 'CUM_QTY' oder 'REAL_QTY'")

    # ✅ Entferne NaN-Werte
    df = df[['CUM_QTY', 'REAL_QTY']].dropna()

    # ✅ **Ausreißerentfernung mit IQR-Methode**
    Q1 = df.quantile(0.25)
    Q3 = df.quantile(0.75)
    IQR = Q3 - Q1
    df_clean = df[~((df < (Q1 - 1.5 * IQR)) | (df > (Q3 + 1.5 * IQR))).any(axis=1)]

    # ✅ Sicherstellen, dass nach der Bereinigung noch genügend Daten vorhanden sind
    if df_clean.shape[0] < 2:
        raise ValueError("Zu wenige Datenpunkte nach der Ausreißer-Entfernung für eine Regression!")

    # ✅ Lineare Regression durchführen
    X = df_clean[['CUM_QTY']].values.reshape(-1, 1)  # Feature
    y = df_clean['REAL_QTY'].values.reshape(-1, 1)   # Zielvariable

    reg = LinearRegression()
    reg.fit(X, y)

    # ✅ Berechnung der Regressionsgeraden
    x_range = np.linspace(df_clean['CUM_QTY'].min(), df_clean['CUM_QTY'].max(), 100).reshape(-1, 1)
    y_pred = reg.predict(x_range)
    file_name = f"Regression_{part}.png"
    # ✅ Scatter-Plot erstellen
    plt.figure(figsize=(12, 6))
    plt.scatter(df_clean['CUM_QTY'], df_clean['REAL_QTY'], color='red', alpha=0.7, label="Realer Bestand vs. Kumulierter Bestand", s=20)
    plt.plot(x_range, y_pred, color='blue', linestyle='-', linewidth=2, label=f"Lineare Regression:\n y = {reg.coef_[0][0]:.4f} * x + {reg.intercept_[0]:.2f}")

    # ✅ Formatierung
    plt.xlabel("Kumulierter Bestand (CUM_QTY)")
    plt.ylabel("Realer Bestand (REAL_QTY)")
    plt.title(f"Scatter-Plot mit linearer Regression: {titletext}")
    plt.legend()
    plt.grid(True, linestyle="--", alpha=0.5)

    # ✅ Diagramm speichern und anzeigen
    # plt.show()
    # ✅ **Grafik speichern**
    plt.savefig(file_name, dpi=300, bbox_inches='tight')
    # ✅ JSON-Datei aktualisieren oder erstellen
    slope_value = round(reg.coef_[0][0], 4)

    # Falls JSON-Datei existiert, lade vorhandene Daten
    if os.path.exists(json_file):
        with open(json_file, "r") as file:
            try:
                data = json.load(file)
            except json.JSONDecodeError:
                data = {}
    else:
        data = {}

    # ✅ Daten für `part` setzen oder aktualisieren
    data[str(part)] = {"slope": slope_value}

    # ✅ JSON-Datei speichern
    with open(json_file, "w") as file:
        json.dump(data, file, indent=4)

    print(f"✅ Regressionsdaten für {part} gespeichert: slope = {slope_value}")

    return df_clean, reg


import json
import os

def get_slope_from_json(part, json_file="regression_results.json"):
    """
    Liest den gespeicherten Slope-Wert für ein bestimmtes Part aus der JSON-Datei aus.

    :param part: ID des Teils als String oder Integer
    :param json_file: Name der JSON-Datei mit den gespeicherten Regressionswerten
    :return: Der Slope-Wert als Float oder None, falls nicht gefunden
    """
    # Falls die Datei nicht existiert, gib None zurück
    if not os.path.exists(json_file):
        print(f"⚠️ Datei '{json_file}' existiert nicht.")
        return None

    # Lade die JSON-Datei
    try:
        with open(json_file, "r") as file:
            data = json.load(file)
    except json.JSONDecodeError:
        print(f"⚠️ Fehler beim Einlesen der Datei '{json_file}'. Die Datei könnte beschädigt sein.")
        return None

    # Suche nach dem angegebenen `part`
    part_str = str(part)  # Sicherstellen, dass der Key als String vorliegt
    if part_str in data:
        slope_value = data[part_str].get("slope", None)
        print(f"✅ Slope für Part {part}: {slope_value}")
        return slope_value
    else:
        print(f"❌ Kein Eintrag für Part {part} gefunden.")
        return None

def process_inventory_trends_with_moving_average(df, grenzdatum_start, grenzdatum_ende, cut_off_date, use_slope = True, calc_slope = False, real_flag = True, initial_stock=None, window=14):
    """
    Berechnet den Bestandsverlauf basierend auf Produktion, Verkauf und Einkauf.
    Zusätzlich wird ein gleitender Mittelwert (Moving Average) für die reale und geschätzte Bestandskurve berechnet.
    
    :param df: Pandas DataFrame mit Bestandsdaten
    :param initial_stock: Optionaler Startbestand, falls keiner aus df['MVMT_ON_HAND'] genutzt werden soll.
    :param window: Fenstergröße für den Moving Average (Standard: 14 Tage)
    :return: DataFrame mit geschätztem Bestand & realem Bestandsverlauf
    """

    # ✅ Sicherstellen, dass alle Datumsfelder als datetime erkannt werden
    date_fields = ['PROD_ORDER_END_DATE', 'ORDER_DOC_LINE_REQUESTED_DATE', 'PUR_ORDER_LINE_REQUESTED_DATE']
    part_id = str(df['PART_I_D'].dropna().iloc[0]) if 'PART_I_D' in df.columns and not df['PART_I_D'].dropna().empty else "UNKNOWN"
    for field in date_fields:
        if field in df.columns:
            df[field] = pd.to_datetime(df[field], errors='coerce')
    # Sicherstellen, dass cut_off_date ebenfalls ein datetime-Objekt ist
    cut_off_date = pd.to_datetime(cut_off_date, errors='coerce')


    # ✅ Produktion extrahieren (negative Menge, da Verbrauch)
    prod_df = df[['PROD_ORDER_END_DATE', 'PROD_ORDER_TARGET_QTY','PROD_ORDER_CREATION_DATE']].dropna()
    # Sicherstellen, dass die Spalte im richtigen Datumsformat ist
    if real_flag:
        prod_df['creation_date'] = pd.to_datetime(prod_df['PROD_ORDER_CREATION_DATE'], errors='coerce')

        # Filtern der Daten
        prod_df = prod_df[prod_df['creation_date'] <= cut_off_date]

    prod_df = prod_df.rename(columns={'PROD_ORDER_END_DATE': 'DATE', 'PROD_ORDER_TARGET_QTY': 'QTY'})
    if use_slope:
        slope = get_slope_from_json(part_id)

        prod_df['QTY'] *= -slope  
    else:
        prod_df['QTY'] *= -1

    print(prod_df['DATE'] )
    # ✅ Verkauf extrahieren (negative Menge, da Lagerabgang)
    order_df = df[['ORDER_DOC_LINE_REQUESTED_DATE', 'ORDER_DOC_LINE_QTY','ORDER_DOC_CREATION_DATE']].dropna()
    
    if real_flag:
        order_df ['creation_date'] = pd.to_datetime(order_df ['ORDER_DOC_CREATION_DATE'], errors='coerce')

        # Filtern der Daten
        order_df  = order_df [order_df ['creation_date'] <= cut_off_date]
    order_df = order_df.rename(columns={'ORDER_DOC_LINE_REQUESTED_DATE': 'DATE', 'ORDER_DOC_LINE_QTY': 'QTY'})
    order_df['QTY'] *= -1  
    print(order_df['DATE'] )
    # ✅ Einkauf extrahieren (positive Menge, da Lagerzugang)
    pur_df = df[['PUR_ORDER_LINE_REQUESTED_DATE', 'PUR_ORDER_LINE_QTY','PUR_ORDER_CREATION_DATE']].dropna()
    if real_flag:
        pur_df ['creation_date'] = pd.to_datetime(pur_df ['PUR_ORDER_CREATION_DATE'], errors='coerce')

        # Filtern der Daten
        pur_df  = pur_df [pur_df ['creation_date'] <= cut_off_date]    
    pur_df = pur_df.rename(columns={'PUR_ORDER_LINE_REQUESTED_DATE': 'DATE', 'PUR_ORDER_LINE_QTY': 'QTY'})
    pur_df['QTY'] *= 1  
    print(pur_df['DATE'] )

    # ✅ Alle Bewegungen zusammenführen
    all_data = pd.concat([prod_df, order_df, pur_df])

    # ✅ **Doppelte Datumswerte aggregieren**
    daily_inventory = all_data.groupby('DATE', as_index=False)['QTY'].sum()

    # ✅ **Lückenlose Zeitreihe erstellen (fehlende Tage auffüllen)**
    #all_dates = pd.date_range(start=daily_inventory['DATE'].min(), end=daily_inventory['DATE'].max(), freq='D')
    all_dates = pd.date_range(start=grenzdatum_start, end=grenzdatum_ende, freq='D')
    daily_inventory = daily_inventory.set_index('DATE').reindex(all_dates.to_numpy(), fill_value=np.nan)
    daily_inventory = daily_inventory.reset_index()
    daily_inventory = daily_inventory.rename(columns={'index': 'DATE'})
    
    # ✅ **Fehlende Werte mit 0 füllen**
    daily_inventory['QTY'] = daily_inventory['QTY'].fillna(0)
    daily_inventory['CUM_QTY'] =  daily_inventory['QTY'].cumsum()
    # ✅ **Realen Bestand mit Moving Average berechnen**
    real_inventory = df[['MVMT_CREATION_DATE_TIME', 'MVMT_ON_HAND']].dropna()
    real_inventory = real_inventory.rename(columns={'MVMT_CREATION_DATE_TIME': 'DATE', 'MVMT_ON_HAND': 'REAL_QTY'})
    real_inventory = real_inventory.groupby('DATE', as_index=False).last()
    real_inventory = real_inventory.set_index('DATE')
    real_inventory = real_inventory.reindex(all_dates.to_numpy(), fill_value=np.nan, method='ffill')  # Fehlende Werte mit letztem Wert füllen
    real_inventory = real_inventory.reset_index()
    real_inventory = real_inventory.rename(columns={'index': 'DATE'})
    
    # ✅ **Anfangsbestand bestimmen**
    if initial_stock is None and not df.empty:
        # Den letzten bekannten Bestand vor `cut_off_date` nehmen
        real_qty = real_inventory.loc[real_inventory['DATE'] <= cut_off_date, 'REAL_QTY'].dropna().values[-1:]
        daily_qty = daily_inventory.loc[daily_inventory['DATE'] <= cut_off_date, 'CUM_QTY'].dropna().values[-1:]

        # Debugging
        print("Letzter bekannter Real Qty:", real_qty)
        print("Letzter bekannter Daily Qty:", daily_qty)

        # Falls beide Werte fehlen, initialisiere mit 0
        if len(real_qty) == 0 or len(daily_qty) == 0:
            initial_stock = 0
        else:
            initial_stock = real_qty[0] - daily_qty[0]
        

    # ✅ **Kumulierte Bestandsberechnung**
    daily_inventory['CUM_QTY'] =  daily_inventory['CUM_QTY']  + initial_stock

    # ✅ **Vor cut_off_date den geschätzten Bestand an den realen Bestand angleichen**
    daily_inventory.loc[daily_inventory['DATE'] <= cut_off_date, 'CUM_QTY'] = real_inventory.loc[real_inventory['DATE'] <= cut_off_date, 'REAL_QTY'].values

    # ✅ **Moving Average für geschätzte Bestandskurve**
    daily_inventory['CUM_QTY_MA'] = daily_inventory['CUM_QTY'].rolling(window=window, min_periods=1).mean()


    # ✅ **Moving Average für realen Bestand**
    real_inventory['REAL_QTY_MA'] = real_inventory['REAL_QTY'].rolling(window=window, min_periods=1).mean()

    # ✅ **Bestandskurven kombinieren**
    final_inventory = daily_inventory.merge(real_inventory, on= 'DATE', how='outer')


# ✅ **Werte für Titel generieren**
    title_fields = ["PART_I_D", "PART_DESC1", "PART_DESC2", "PART_DESC3", "PART_DESC4"]
    title_values = [str(df[field].dropna().iloc[0]) if field in df.columns and not df[field].dropna().empty else "N/A" for field in title_fields]
    title_text = " - ".join(title_values)
    # ✅ **Dateiname für die Grafik setzen**
    cut_off_date_str = str(cut_off_date)[:10].replace("-", "_")
    
    file_name = f"Bestandsverlauf_{part_id}_cutoffdate_{cut_off_date_str}_realflag_{str(real_flag)}_useslope_{str(use_slope)}.png"
    cutoff_date = pd.to_datetime(cut_off_date)
    end_date = cutoff_date + pd.DateOffset(months=4)
    # ✅ **Plotten der Bestandskurven**
    plt.figure(figsize=(12, 6))
    if real_flag == False and use_slope == False:
        plot_scatter_with_regression(final_inventory, title_text, part_id )
    # Geschätzter Bestand
    plt.plot(final_inventory['DATE'], final_inventory['CUM_QTY'], linestyle='-', color='b', label=f"Geschätzter Bestand am {cut_off_date_str}")

    # Gleitender Mittelwert für geschätzten Bestand
    plt.plot(final_inventory['DATE'], final_inventory['CUM_QTY_MA'], linestyle='-', color='orange', linewidth=2, label=f'{window}-Tage MA Geschätzter Bestand')

    # Reale Bestandswerte
    plt.plot(final_inventory['DATE'], final_inventory['REAL_QTY'], linestyle='--', color='red', alpha=0.5, label="Realer Bestand")

    # Gleitender Mittelwert des realen Bestands
    plt.plot(final_inventory['DATE'], final_inventory['REAL_QTY_MA'], linestyle='-', color='green', linewidth=2, label=f'{window}-Tage MA Realer Bestand')
    # ✅ Schraffierten Bereich hinzufügen (grauer Hintergrund)
    plt.axvspan(cut_off_date, end_date, color='gray', alpha=0.2)
    # **WICHTIG**: X-Achse begrenzen, damit nicht zu viel angezeigt wird
    plt.xlim(pd.to_datetime(grenzdatum_start), pd.to_datetime(grenzdatum_ende))
    # **Diagramm Formatierung**
    plt.gca().xaxis.set_major_formatter(mdates.DateFormatter('%Y-%m-%d'))
    plt.gca().xaxis.set_major_locator(mdates.MonthLocator())
    plt.xticks(rotation=45)

    plt.xlabel("Datum")
    plt.ylabel("Bestand")
    plt.title(f"Bestandsverlauf: {title_text}")
    plt.legend()
    plt.grid()


    # ✅ **Grafik speichern**
    plt.savefig(file_name, dpi=300, bbox_inches='tight')
    print(f"🔹 Grafik gespeichert als: {file_name}")
    

    return final_inventory


# ✅ **Beispielaufruf**
# final_inventory_df = process_inventory_trends_with_moving_average(df, initial_stock=50000, window=14)


# ✅ Beispielaufruf mit df
# daily_inventory_df = process_inventory_trends(df)

# **Beispielaufruf**
# df_moving_avg = bestandsverlauf_moving_average(df_gesamtbestand, window=14)

def process_combinations(use_slope = False,  real_flag=True):
    conn = dbapi.connect(
    address="10.0.0.80",
    port=30015,
    user="CUSTOMER_ADMIN",
    password="mIBFtYdgmE4nh0VaNkJS",
    )
    columns=[
    'MVMT_STORAGE_AREA',
    'MVMT_STORAGE_AREA_ON_HAND',
    'MVMT_CREATION_DATE_TIME',
    'MVMT_ON_HAND',
    'PROD_ORDER_CREATION_DATE',
    'PROD_ORDER_END_DATE',
    'PROD_ORDER_TARGET_QTY',
    'ORDER_DOC_LINE_QTY', 
    'ORDER_DOC_LINE_REQUESTED_DATE',
    'ORDER_DOC_CREATION_DATE',
    'PUR_ORDER_CREATION_DATE',
    'PUR_ORDER_LINE_QTY',
    'PUR_ORDER_LINE_REQUESTED_DATE',
    "PART_DESC1", 
    "PART_DESC2",
    "PART_DESC3",
    "PART_DESC4",
    "PART_I_D",
    ]

    part_list = [
        '31612046',
        '031800112',
        '0240121',
        '044000131',
        '0630015',
        '080205401',
        '0239205',
        '0241030',
        '031800010',
        '0802009',
        '0803026',
        '080206140',
        '0631678',
        '0239204',
        '0803090',
        '0801006',
        '202074070',
        '080206120',
        '031801212',
        '024599490',
        '0630256'
        ]
    
    column_list = ", ".join(f'"{col}"' for col in columns)
    grenzdatum_start = date(2023, 1, 1)  # Python DATE-Objekt
    grenzdatum_ende = date(2024, 12, 31)
    cut_off_date = date(2023, 8, 1)

    for part in part_list:

        query = f"""
            SELECT {column_list} 
            FROM "NEMO"."pa_export_apra" 
            WHERE "COMPANY" = ?
            AND TO_DATE("PROCESS_DATE", 'YYYY-MM-DD') > ?
            AND TO_DATE("PROCESS_DATE", 'YYYY-MM-DD') < ?
            AND (PROCESS = 'Make'
            OR PROCESS = 'Production'
            OR PROCESS = 'Source'
            OR PROCESS = 'Deliver'
            )
            AND PART_I_D = '{part}'
        ;
        """

        params = ("001", grenzdatum_start.strftime("%Y-%m-%d"), grenzdatum_ende.strftime("%Y-%m-%d")) #,cut_off_date.strftime("%Y-%m-%d"),cut_off_date.strftime("%Y-%m-%d"),cut_off_date.strftime("%Y-%m-%d"), grenzdatum_ende.strftime("%Y-%m-%d"))  # Übergabe als String
        results = execute_query(conn, query,params=params)
        
        # Erstellen eines DataFrame aus den Ergebnissen
        df = pd.DataFrame(results, columns=columns)
        if df.empty:
            print("Dataframe ist leer")
            continue
        # df_bestand = bestandsverlauf_pro_LO_tag(df)
        #final_inventory_df = process_inventory_trends_with_moving_average(df,grenzdatum_start,grenzdatum_ende, cut_off_date.strftime("%Y-%m-%d"), real_flag=True , initial_stock=None, window=34)
        try:
            
            final_inventory_df = process_inventory_trends_with_moving_average(df,grenzdatum_start,grenzdatum_ende, cut_off_date.strftime("%Y-%m-%d"), use_slope = use_slope, calc_slope=False, real_flag=real_flag , initial_stock=None, window=34)
            #df_gesamtbestand = gesamter_bestandsverlauf(df)
            #df_moving_avg = bestandsverlauf_moving_average(df_gesamtbestand, window=14)
            #final_inventory_df = process_inventory_trends_with_moving_average(df, cut_off_date, initial_stock=None, window=34)
        except Exception as e:
            print('Fehler in: ' + str(part) + ' ' + str(e))  
        #break 
    conn.close()
    #plt.show()
    #'PROD_ORDER_END_DATE', 'PROD_ORDER_TARGET_QTY',
    # 'ORDER_DOC_LINE_QTY', 'ORDER_DOC_LINE_REQUESTED_DATE',
    # 'PUR_ORDER_LINE_QTY', 'PUR_ORDER_LINE_REQUESTED_DATE',  
def create_pptx_from_images(image_folder, output_pptx="Bestandsverlauf_Presentation.pptx"):
    """
    Erstellt eine PowerPoint-Präsentation mit den generierten Bestandsverlauf-Bildern.
    
    - Jedes Part_ID bekommt eine eigene Folie.
    - Pro Folie werden 4 Bilder in einem festen Raster platziert.
    - Die Reihenfolge der Positionen ist immer:
      1. (False, False) → **Unten links**
      2. (True, False) → **Oben links**
      3. (True, True) → **Oben rechts**
      4. (False, True) → **Unten rechts**
    
    :param image_folder: Verzeichnis, in dem die Bilder gespeichert sind.
    :param output_pptx: Name der generierten PowerPoint-Datei.
    """

    # PowerPoint-Präsentation initialisieren
    prs = Presentation()

    # Regex-Muster zum Parsen der Dateinamen
    image_pattern = re.compile(r"Bestandsverlauf_(.*?)_cutoffdate_(.*?)_realflag_(.*?)_useslope_(.*?)\.png")
    regression_pattern = re.compile(r"Regression_(.*?).png")  # Regex für Regressionsbilder
    images_by_part = {}
    regression_images = {}

    
    # Bilder gruppieren nach Part-ID und Cutoff-Date
    for file in os.listdir(image_folder):
        match = image_pattern.match(file)
        if match:
            part_id, cutoff_date, real_flag, use_slope = match.groups()
            key = f"{part_id}_{cutoff_date}"
            if part_id not in images_by_part:
                images_by_part[part_id] = []
            images_by_part[part_id].append((file, real_flag, use_slope))
            # Regressionsbilder einlesen
        reg_match = regression_pattern.match(file)
        if reg_match:
            part_id  = reg_match.groups()
            regression_images[f"{part_id[0]}"] = file
    print(images_by_part)
    print(regression_images)
    # Definierte Positionen für die vier Varianten
    fixed_positions = {
        ("False", "False"): (Inches(0.5), Inches(4.5)),  # 🔹 (1) Unten links
        ("True", "False"): (Inches(0.5), Inches(1.5)),   # 🔹 (2) Oben links
        ("True", "True"): (Inches(5.5), Inches(1.5)),    # 🔹 (3) Oben rechts
        ("False", "True"): (Inches(5.5), Inches(4.5))    # 🔹 (4) Unten rechts
    }
    regression_position = (Inches(3), Inches(3))  # 📍 Zentrum der Folie für Regression
    # Schleife über alle Part-IDs und Folien erstellen
    for part_key, images in images_by_part.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Leere Folie

        # Bilder platzieren
        for file, real_flag, use_slope in images:
            img_path = os.path.join(image_folder, file)
            position = (real_flag, use_slope)

            if position in fixed_positions:
                x, y = fixed_positions[position]
                slide.shapes.add_picture(img_path, x, y, width=Inches(4))

                # Beschriftung für real_flag und use_slope
                text_box = slide.shapes.add_textbox(x, y - Inches(0.4), Inches(4), Inches(0.4))
                text_frame = text_box.text_frame
                text_frame.text = f"real_flag: {real_flag}, use_slope: {use_slope}"
        # 📌 Regressionsbild einfügen (falls vorhanden)
        print(part_key)
        if part_key in regression_images:
            reg_img_path = os.path.join(image_folder, regression_images[part_key])
            slide.shapes.add_picture(reg_img_path, *regression_position, width=Inches(4))

    # PowerPoint speichern
    prs.save(output_pptx)
    print(f"✅ PowerPoint gespeichert als {output_pptx}")


if __name__ == '__main__':
    #process_combinations(use_slope=True,real_flag=True)
    #process_combinations(use_slope=True,real_flag=False)
    #process_combinations(use_slope=False,real_flag=True)
    #process_combinations(use_slope=False,real_flag=False)
    create_pptx_from_images(os.getcwd(), "Bestandsverlaeufe.pptx")